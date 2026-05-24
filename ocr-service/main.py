from __future__ import annotations

import io
import json
import os

from dotenv import load_dotenv

import re
from html.parser import HTMLParser
from dataclasses import dataclass
from typing import Callable


import cv2
import numpy as np
import pytesseract
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from PIL import Image, UnidentifiedImageError

try:
    import fitz
except ImportError:
    fitz = None

try:
    from google import genai
    from google.genai import types
except ImportError:
    genai = None
    types = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

load_dotenv()


app = FastAPI()


@dataclass
class ExtractedField:
    value: str
    confidence: float
    source: str


FIELD_ALIASES = {
    "first_name": ("first name", "given name", "student first name"),
    "middle_name": ("middle name", "second name", "middle initial"),
    "last_name": ("last name", "surname", "family name", "student last name"),
    "name": ("name", "full name", "student name", "applicant name", "candidate name"),
    "father_name": ("father name", "father's name", "father full name", "parent name", "parent's name"),
    "gender": ("gender", "sex"),
    "date_of_birth": ("date of birth", "dob", "d o b", "d.o.b", "birth date", "birthday"),
    "address_line_1": ("address line 1", "address line one", "address 1"),
    "address_line_2": ("address line 2", "address line two", "address 2"),
    "address": ("address", "residential address", "permanent address", "current address", "mailing address"),
    "city": ("city", "town"),
    "state": ("state", "province"),
    "pin_code": ("pin code", "pincode", "postal code", "zip code"),
    "phone": ("phone", "mobile", "mobile number", "phone number", "contact", "contact number", "cell"),
    "email": ("email", "email id", "e-mail", "e-mail id", "mail"),
    "branch": ("branch", "department", "course", "programme", "program", "stream", "academic branch"),
    "student_id": ("student id", "student identity", "roll number", "roll no", "registration number", "admission number", "enrollment number", "application number", "id number", "id"),
}

GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
GEMINI_EXTRACTION_PROMPT = """
Analyze the uploaded document intelligently, not just as OCR text extraction.

Your task is to fully understand the document visually and semantically.

Instructions:

Extract ALL visible text from the document.
Understand what the document is.
Identify:
document type
organization/institution name
headings
subheadings
labels
values
IDs
years
signatures
addresses
phone numbers
dates
logos
stamps
side text
vertical text
footer text
handwritten text
metadata
codes without labels
important standalone text
Do NOT limit extraction to only label-value pairs.
Detect information even if no label exists nearby.
Preserve the actual meaning and hierarchy of the document.
Understand contextual relationships between texts.
"""

GEMINI_RESPONSE_SCHEMA = {
    "type": "object",
    "properties": {
        "fields": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "label": {"type": "string"},
                    "key": {"type": "string"},
                    "value": {"type": "string"},
                    "confidence": {"type": "number"},
                },
                "required": ["label", "key", "value", "confidence"],
                "additionalProperties": False,
            },
        }
    },
    "required": ["fields"],
    "additionalProperties": False,
}

GEMINI_PAGE_CONTEXT_PROMPT = """
Analyze the screenshot of a webpage containing a form or input fields.
Return JSON only. Identify visible field labels that correspond to the following canonical keys:
first_name, middle_name, last_name, name, father_name, gender, date_of_birth,
address_line_1, address_line_2, address, city, state, pin_code, phone, email,
branch, student_id.

Rules:
- Output exactly one JSON object.
- Use property `field_labels` as an array of objects.
- Each object must contain `key` and `label`.
- `key` must be one of the canonical keys above.
- `label` must be the exact visible label text shown on the page.
- Do not guess values, only return labels visible around page inputs.
- If a label is repeated for the same key, include it once.
"""

GEMINI_FIELD_MAPPING_PROMPT = """
You are a smart AI assistant that maps extracted document fields to webpage form inputs.
You are given a screenshot of the webpage and a list of extracted keys from a document.
You also have a list of candidate DOM controls with their visible labels, placeholders, ids, names, and nearby text.
Your job is to map each extracted key to the best matching webpage control based on meaning, not exact wording.

Rules:
- Use semantic meaning of the key and the extracted label/value to choose the right page field.
- Prefer controls whose visible label text or nearby prompt text expresses the same concept.
- If the key refers to the same thing as a form label, even if wording differs, map it.
- If the control is clearly the same field but uses a different label style or phrasing, that is okay.
- Use `id` if present, otherwise `name`, otherwise `xpath`.
- Only map keys to actual controls on the page; do not invent fields.
- If you are not confident, set selector:null and confidence:0.
- Return JSON only.

Expected output:
{"mappings":[{"key":"...","selector":{"by":"id"|"name"|"xpath","value":"..."},"confidence":0.0}]}
"""

GEMINI_FIELD_MAPPING_SCHEMA = {
    "type": "object",
    "properties": {
        "mappings": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "key": {"type": "string"},
                    "selector": {
                        "type": ["object", "null"],
                        "properties": {
                            "by": {"type": "string", "enum": ["id", "name", "xpath"]},
                            "value": {"type": "string"},
                        },
                        "required": ["by", "value"],
                        "additionalProperties": False,
                    },
                    "confidence": {"type": "number"},
                },
                "required": ["key", "selector", "confidence"],
                "additionalProperties": False,
            },
        }
    },
    "required": ["mappings"],
    "additionalProperties": False,
}

GEMINI_PAGE_CONTEXT_SCHEMA = {
    "type": "object",
    "properties": {
        "field_labels": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "key": {"type": "string"},
                    "label": {"type": "string"},
                },
                "required": ["key", "label"],
                "additionalProperties": False,
            },
        }
    },
    "required": ["field_labels"],
    "additionalProperties": False,
}

TEXT_EXTENSIONS = {
    ".txt",
    ".csv",
    ".tsv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".md",
    ".log",
}

IMAGE_EXTENSIONS = {
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
    ".tif",
    ".tiff",
}


def pil_to_cv2(image: Image.Image) -> np.ndarray:
    return cv2.cvtColor(np.array(image.convert("RGB")), cv2.COLOR_RGB2BGR)


def preprocess_image(image: Image.Image) -> Image.Image:
    cv_image = pil_to_cv2(image)
    gray = cv2.cvtColor(cv_image, cv2.COLOR_BGR2GRAY)

    scale = max(1.0, 1500 / max(gray.shape[:2]))
    if scale > 1:
        gray = cv2.resize(gray, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)

    gray = cv2.fastNlMeansDenoising(gray, h=20)
    gray = cv2.convertScaleAbs(gray, alpha=1.25, beta=10)
    binary = cv2.adaptiveThreshold(
        gray,
        255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        31,
        11,
    )

    return Image.fromarray(binary)


def image_to_png_bytes(image: Image.Image) -> bytes:
    output = io.BytesIO()
    image.convert("RGB").save(output, format="PNG")
    return output.getvalue()


def ocr_image(image: Image.Image) -> tuple[str, float]:
    processed = preprocess_image(image)
    text = pytesseract.image_to_string(processed, config="--oem 3 --psm 6")
    data = pytesseract.image_to_data(processed, output_type=pytesseract.Output.DICT)

    confidences = []
    for value, word in zip(data.get("conf", []), data.get("text", [])):
        try:
            confidence = float(value)
        except (TypeError, ValueError):
            continue

        if confidence >= 0 and word.strip():
            confidences.append(confidence)

    average_confidence = round(sum(confidences) / len(confidences), 2) if confidences else 0.0
    return text, average_confidence


def extract_page_hint_labels(text: str) -> list[str]:
    lines = [normalize_text(line) for line in text.splitlines() if normalize_text(line)]
    labels: list[str] = []

    for line in lines:
        if len(line) < 3 or len(line) > 120:
            continue

        parts = re.split(r"\s*(?:[:=\|\-])\s*", line, maxsplit=1)
        label_candidate = parts[0].strip()
        if 2 <= len(label_candidate) <= 60:
            labels.append(label_candidate)
            continue

        if len(line.split()) <= 5:
            labels.append(line)

    return list(dict.fromkeys(labels))


def extract_page_label_map(text: str) -> dict[str, list[str]]:
    labels = extract_page_hint_labels(text)
    label_map: dict[str, list[str]] = {}

    for label in labels:
        key = canonical_key_for_label(label)
        if not key:
            continue

        label_map.setdefault(key, []).append(label)

    return {key: list(dict.fromkeys(values)) for key, values in label_map.items()}


def response_from_pages(page_results: list[dict], page_context: dict | None = None) -> dict:
    merged_fields = merge_fields(page_results)

    response = {
        "raw_text": "\n\n".join(page["raw_text"] for page in page_results),
        "ocr_confidence": round(
            sum(page["ocr_confidence"] for page in page_results) / len(page_results),
            2,
        ),
        "data": {key: field.value for key, field in merged_fields.items()},
        "fields": {
            key: {
                "value": field.value,
                "confidence": field.confidence,
                "source": field.source,
            }
            for key, field in merged_fields.items()
        },
        "extraction_method": "gemini" if any(page["extraction_method"] == "gemini" for page in page_results) else page_results[0]["extraction_method"],
        "gemini_status": next((page["gemini_status"] for page in page_results if page["gemini_status"] == "used"), page_results[0]["gemini_status"]),
        "pages": [
            {
                "page": page["page"],
                "ocr_confidence": page["ocr_confidence"],
                "extraction_method": page["extraction_method"],
                "gemini_status": page["gemini_status"],
                "data": {key: field.value for key, field in page["fields"].items()},
            }
            for page in page_results
        ],
    }

    if page_context is not None:
        response["page_context"] = page_context

    return response


def gemini_client() -> tuple[object | None, str]:
    if genai is None or types is None:
        return None, "google-genai package is not installed"

    if not os.getenv("GEMINI_API_KEY") and not os.getenv("GOOGLE_API_KEY"):
        return None, "GEMINI_API_KEY is not set"

    return genai.Client(), "ready"


def fields_from_gemini_payload(payload: dict) -> dict[str, ExtractedField]:
    fields = {}

    for item in payload.get("fields", []):
        if not isinstance(item, dict):
            continue

        label = clean_plain_text(item.get("label", ""))
        key = normalize_key(item.get("key") or label)
        if not key:
            continue
        if should_skip_extracted_field(label, item.get("value", "")):
            continue

        canonical_key = canonical_key_for_label(f"{label} {key}")
        output_key = canonical_key or key
        value = clean_field_value(output_key, item.get("value", ""))
        if not value:
            continue

        confidence = item.get("confidence", 0)
        try:
            confidence = float(confidence)
        except (TypeError, ValueError):
            confidence = 0.0

        field = ExtractedField(
            value=value,
            confidence=round(max(0.0, min(100.0, confidence)), 2),
            source=label or "gemini",
        )

        if output_key not in fields or field.confidence > fields[output_key].confidence:
            fields[output_key] = field

    return fields


def extract_fields_with_gemini(image: Image.Image) -> tuple[dict[str, ExtractedField], str]:
    client, status = gemini_client()
    if client is None:
        return {}, "skipped"

    try:
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[
                types.Part.from_bytes(
                    data=image_to_png_bytes(image),
                    mime_type="image/png",
                ),
                GEMINI_EXTRACTION_PROMPT,
            ],
            config={
                "response_mime_type": "application/json",
                "response_json_schema": GEMINI_RESPONSE_SCHEMA,
            },
        )
    except Exception as exc:
        return {}, "fallback"

    try:
        payload = json.loads(response.text)
    except json.JSONDecodeError:
        return {}, "fallback"

    fields = fields_from_gemini_payload(payload)

    if not fields:
        return {}, "fallback"

    return fields, "used"


def extract_text_fields_with_gemini(text: str) -> tuple[dict[str, ExtractedField], str]:
    client, status = gemini_client()
    if client is None:
        return {}, "skipped"

    try:
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[
                GEMINI_EXTRACTION_PROMPT,
                "\nDocument text:\n",
                text[:12000],
            ],
            config={
                "response_mime_type": "application/json",
                "response_json_schema": GEMINI_RESPONSE_SCHEMA,
            },
        )
    except Exception as exc:
        return {}, "fallback"

    try:
        payload = json.loads(response.text)
    except json.JSONDecodeError:
        return {}, "fallback"

    fields = fields_from_gemini_payload(payload)
    if not fields:
        return {}, "fallback"

    return fields, "used"


def images_from_pdf(pdf_bytes: bytes) -> list[Image.Image]:
    if fitz is None:
        raise HTTPException(
            status_code=400,
            detail="PDF support requires PyMuPDF. Install it with: pip install PyMuPDF",
        )

    document = fitz.open(stream=pdf_bytes, filetype="pdf")
    images = []

    for page in document:
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        image = Image.open(io.BytesIO(pixmap.tobytes("png")))
        images.append(image)

    return images


def images_from_upload(file_bytes: bytes, content_type: str | None, filename: str | None) -> list[Image.Image]:
    is_pdf = content_type == "application/pdf" or (filename or "").lower().endswith(".pdf")
    if is_pdf:
        return images_from_pdf(file_bytes)

    try:
        return [Image.open(io.BytesIO(file_bytes))]
    except UnidentifiedImageError as exc:
        raise HTTPException(status_code=400, detail="Upload must be an image or PDF file.") from exc


class TextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []

    def handle_data(self, data: str):
        if data.strip():
            self.parts.append(data.strip())

    def text(self) -> str:
        return "\n".join(self.parts)


def file_extension(filename: str | None) -> str:
    _, extension = os.path.splitext(filename or "")
    return extension.lower()


def decode_text(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def extract_text_from_plain_file(file_bytes: bytes, extension: str) -> str:
    text = decode_text(file_bytes)
    if extension in {".html", ".htm", ".xml"}:
        parser = TextExtractor()
        parser.feed(text)
        return parser.text()
    return text


def extract_text_from_docx(file_bytes: bytes) -> str:
    if docx is None:
        raise HTTPException(status_code=400, detail="DOCX support requires python-docx.")

    document = docx.Document(io.BytesIO(file_bytes))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]

    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                parts.append(" | ".join(values))

    return "\n".join(parts)


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    if openpyxl is None:
        raise HTTPException(status_code=400, detail="XLSX support requires openpyxl.")

    workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []

    for sheet in workbook.worksheets:
        parts.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                parts.append(" | ".join(values))

    return "\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    if Presentation is None:
        raise HTTPException(status_code=400, detail="PPTX support requires python-pptx.")

    presentation = Presentation(io.BytesIO(file_bytes))
    parts = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"Slide {slide_number}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())

    return "\n".join(parts)


def text_from_upload(file_bytes: bytes, filename: str | None) -> tuple[str, str] | None:
    extension = file_extension(filename)

    if extension in TEXT_EXTENSIONS:
        return extract_text_from_plain_file(file_bytes, extension), "text-file"
    if extension == ".docx":
        return extract_text_from_docx(file_bytes), "docx"
    if extension in {".xlsx", ".xlsm"}:
        return extract_text_from_xlsx(file_bytes), "xlsx"
    if extension == ".pptx":
        return extract_text_from_pptx(file_bytes), "pptx"

    return None


def upload_kind(file_bytes: bytes, content_type: str | None, filename: str | None) -> str:
    extension = file_extension(filename)
    if extension == ".pdf" or content_type == "application/pdf":
        return "image-pages"
    if extension in IMAGE_EXTENSIONS or (content_type or "").startswith("image/"):
        return "image-pages"
    if text_from_upload(file_bytes, filename) is not None:
        return "text"
    return "image-pages"


def fields_from_gemini_context_payload(payload: dict) -> dict[str, list[str]]:
    field_map: dict[str, list[str]] = {}

    for item in payload.get("field_labels", []):
        if not isinstance(item, dict):
            continue

        key = normalize_key(item.get("key") or "")
        label = normalize_text(item.get("label") or "")
        if not key or not label:
            continue

        if key not in FIELD_ALIASES:
            continue

        field_map.setdefault(key, []).append(label)

    return {key: list(dict.fromkeys(labels)) for key, labels in field_map.items()}


def extract_field_mapping_with_gemini(
    image: Image.Image,
    extracted: dict[str, object],
    dom_controls: list[dict],
    field_sources: dict[str, str] | None = None,
    extracted_text: str | None = None,
    source_filename: str | None = None,
) -> tuple[list[dict], str]:
    client, status = gemini_client()
    if client is None:
        return [], "skipped"

    control_lines = []
    for control in dom_controls:
        parts = [control.get('labelText') or '', control.get('nearestText') or '', control.get('placeholder') or '', control.get('exactText') or '']
        label = normalize_text(' | '.join([p for p in parts if p]))
        control_lines.append(
            f"id={control.get('id') or ''}; name={control.get('name') or ''}; type={control.get('type') or ''}; label={label}; xpath={control.get('xpath') or ''}"
        )

    extracted_lines = []
    for key, value in extracted.items():
        source_label = field_sources.get(key) if field_sources else None
        line = f"- {key}: {normalize_text(str(value))}"
        if source_label:
            line += f" (label: {normalize_text(source_label)})"
        extracted_lines.append(line)

    document_hint = ''
    if source_filename:
        document_hint = f"The source document is attached as {source_filename}."

    if extracted_text:
        extracted_text = normalize_text(extracted_text)
        if len(extracted_text) > 4000:
            extracted_text = extracted_text[:4000]

    prompt_parts = [
        "You are mapping extracted document fields to a webpage form.",
        document_hint,
        "Use the meaning of each source field label and value to select the best matching input control.",
        "If the form question text is semantically equivalent, map it even when wording differs.",
        "Do not guess fields that are not present in the form.",
        "Return only valid control selectors with by=id, name, or xpath.",
        "Target keys include: name, father_name, date_of_birth, branch, student_id, address, phone, email, gender, city, state, pin_code.",
        "Examples:",
        "- If document field label is 'Name' and form question says 'Student's full name', map to that field.",
        "- If document field label is 'F. Name' and form says 'Father's name', map to father_name.",
        "- If document field label is 'Branch' and form says 'Course / Branch', map to branch.",
        "- If document field label is 'Roll Number' and form says 'Student ID', map to student_id.",
        "Extracted keys and values:",
        *extracted_lines,
    ]

    if extracted_text:
        prompt_parts.extend(["Document extracted text:", extracted_text])

    prompt_parts.extend([
        "Web form candidate controls:",
        *control_lines,
        "Return JSON only with this format:",
        '{"mappings":[{"key":"...","selector":{"by":"id"|"name"|"xpath","value":"..."},"confidence":0.0}]}'
    ])

    prompt = "\n".join([part for part in prompt_parts if part])

    try:
        contents = [types.Part.from_bytes(data=image_to_png_bytes(image), mime_type="image/png")]
        contents.append(prompt)

        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=contents,
            config={
                "response_mime_type": "application/json",
                "response_json_schema": GEMINI_FIELD_MAPPING_SCHEMA,
            },
        )
    except Exception:
        return [], "fallback"

    try:
        payload = json.loads(response.text)
    except json.JSONDecodeError:
        return [], "fallback"

    mappings = []
    for item in payload.get('mappings', []):
        if not isinstance(item, dict):
            continue
        key = item.get('key')
        selector = item.get('selector')
        confidence = item.get('confidence', 0)
        if not key or selector is None:
            mappings.append({"key": key or "", "selector": None, "confidence": 0})
            continue
        if not isinstance(selector, dict):
            mappings.append({"key": key, "selector": None, "confidence": 0})
            continue
        by = selector.get('by')
        value = selector.get('value')
        if by not in {'id', 'name', 'xpath'} or not value:
            mappings.append({"key": key, "selector": None, "confidence": 0})
            continue
        mappings.append({"key": key, "selector": {"by": by, "value": value}, "confidence": float(confidence)})

    return mappings, "used"


def extract_page_context_with_gemini(image: Image.Image) -> tuple[dict[str, list[str]], str]:
    client, status = gemini_client()
    if client is None:
        return {}, "skipped"

    try:
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[
                types.Part.from_bytes(data=image_to_png_bytes(image), mime_type="image/png"),
                GEMINI_PAGE_CONTEXT_PROMPT,
            ],
            config={
                "response_mime_type": "application/json",
                "response_json_schema": GEMINI_PAGE_CONTEXT_SCHEMA,
            },
        )
    except Exception:
        return {}, "fallback"

    try:
        payload = json.loads(response.text)
    except json.JSONDecodeError:
        return {}, "fallback"

    page_map = fields_from_gemini_context_payload(payload)
    if not page_map:
        return {}, "fallback"

    return page_map, "used"


def build_page_context_from_screenshot(image: Image.Image) -> dict:
    text, confidence = ocr_image(image)
    labels = extract_page_hint_labels(text)
    label_map = extract_page_label_map(text)
    ai_label_map, ai_status = extract_page_context_with_gemini(image)

    merged_label_map = {**label_map}
    for key, items in ai_label_map.items():
        merged_label_map.setdefault(key, [])
        merged_label_map[key].extend(items)
        merged_label_map[key] = list(dict.fromkeys(merged_label_map[key]))

    page_context = {
        "screenshot_text": normalize_text(text)[:4000],
        "screenshot_ocr_confidence": confidence,
        "hint_labels": labels,
        "field_label_map": merged_label_map,
        "ai_page_context_status": ai_status or "skipped",
    }

    if ai_label_map:
        page_context["ai_field_label_map"] = ai_label_map
    return page_context


def normalize_text(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def normalize_key(value: object) -> str:
    key = re.sub(r"[^a-z0-9]+", "_", str(value).lower()).strip("_")
    return re.sub(r"_+", "_", key)


def should_skip_extracted_field(label: object, value: object) -> bool:
    label_text = normalize_text(str(label)).lower()
    value_text = normalize_text(str(value)).lower()
    combined = f"{label_text} {value_text}"

    noise_phrases = (
        "college of engineering",
        "student identity card",
        "nba accredited",
        "aicte approved",
        "affiliated",
        "autonomous institution",
    )

    if any(phrase in combined for phrase in noise_phrases):
        return True

    if "signature" in label_text:
        return True

    return False


def canonical_key_for_label(label: str) -> str | None:
    label = normalize_text(re.sub(r"[^A-Za-z0-9 ]", " ", label)).lower()
    best_key = None
    best_score = 0.0

    for key, aliases in FIELD_ALIASES.items():
        for alias in aliases:
            score = similarity(label, alias)
            if alias in label:
                score = max(score, 1.0)
            if score > best_score:
                best_key = key
                best_score = score

    return best_key if best_score >= 0.76 else None


def clean_field_value(key: str, value: object) -> str:
    if key in {"first_name", "middle_name", "last_name", "name"}:
        return clean_name(str(value))
    if key == "phone":
        return clean_phone(str(value))
    if key == "email":
        return clean_email(str(value))
    if key == "pin_code":
        return clean_pin_code(str(value))
    return clean_plain_text(str(value))


def clean_name(value: str) -> str:
    value = re.split(r"\b(?:phone|mobile|email|e-mail|dob|address)\b", value, flags=re.IGNORECASE)[0]
    return normalize_text(re.sub(r"[^A-Za-z .'-]", " ", value))


def clean_plain_text(value: str) -> str:
    return normalize_text(str(value))


def clean_pin_code(value: str) -> str:
    match = re.search(r"\b\d{5,6}\b", str(value))
    return match.group(0) if match else normalize_text(str(value))


def clean_phone(value: str) -> str:
    match = re.search(r"(?:\+?\d[\d\s().-]{7,}\d)", value)
    if not match:
        return ""

    phone = normalize_text(match.group(0))
    digits = re.sub(r"\D", "", phone)
    looks_like_date = bool(re.fullmatch(r"\d{1,2}[-/.]\d{1,2}[-/.]\d{2,4}", phone))

    if looks_like_date or len(digits) < 10 or len(digits) > 15:
        return ""

    return phone


def clean_email(value: str) -> str:
    value = re.sub(r"\s*@\s*", "@", value)
    value = re.sub(r"\s*\.\s*", ".", value)
    match = re.search(r"[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}", value, re.IGNORECASE)
    return match.group(0).strip() if match else ""


def find_first_email(text: str) -> str:
    match = re.search(r"[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}", text, re.IGNORECASE)
    return match.group(0).strip() if match else ""


def find_first_phone(text: str) -> str:
    for match in re.finditer(r"(?:\+?\d[\d\s().-]{7,}\d)", text):
        candidate = normalize_text(match.group(0))
        digits = re.sub(r"\D", "", candidate)
        if 10 <= len(digits) <= 15 and not re.fullmatch(r"\d{1,2}[-/.]\d{1,2}[-/.]\d{2,4}", candidate):
            return candidate
    return ""


def guess_name_from_text(text: str, email: str, phone: str) -> str:
    lines = [normalize_text(line) for line in text.splitlines() if normalize_text(line)]
    common_sections = re.compile(r"\b(summary|objective|skills|projects|experience|education|achievements|certifications|technical|languages|tools|contact|resume|curriculum vitae)\b", re.I)

    for line in lines[:6]:
        if email and email.lower() in line.lower():
            continue
        if phone and phone in line:
            continue
        if common_sections.search(line):
            continue
        if re.search(r"\b(linkedin|github|leetcode|portfolio|website)\b", line, re.I):
            continue
        if 2 <= len(line.split()) <= 5 and re.search(r"[A-Za-z]", line):
            return line

    if email:
        local_part = email.split("@", 1)[0]
        name_candidate = re.sub(r"[\d._-]+", " ", local_part).strip()
        if 1 < len(name_candidate) <= 40 and re.search(r"[A-Za-z]", name_candidate):
            return normalize_text(name_candidate)

    return ""


def best_label_match(line: str, aliases: tuple[str, ...]) -> float:
    label = re.split(r"[:|=-]", line, maxsplit=1)[0].lower()
    label = re.sub(r"[^a-z0-9 ]", " ", label)
    label = normalize_text(label)

    if not label:
        return 0.0

    scores = [similarity(label, alias) for alias in aliases]
    scores.extend(1.0 for alias in aliases if alias in label)
    return max(scores, default=0.0)


def similarity(left: str, right: str) -> float:
    left = left.lower()
    right = right.lower()
    if left == right:
        return 1.0

    previous = list(range(len(right) + 1))
    for i, left_char in enumerate(left, start=1):
        current = [i]
        for j, right_char in enumerate(right, start=1):
            insert = current[j - 1] + 1
            delete = previous[j] + 1
            replace = previous[j - 1] + (left_char != right_char)
            current.append(min(insert, delete, replace))
        previous = current

    distance = previous[-1]
    return 1 - (distance / max(len(left), len(right), 1))


def extract_by_label(
    lines: list[str],
    field_name: str,
    cleaner: Callable[[str], str],
    fallback_pattern: str | None = None,
) -> ExtractedField | None:
    aliases = FIELD_ALIASES[field_name]
    candidates: list[ExtractedField] = []

    for index, line in enumerate(lines):
        score = best_label_match(line, aliases)
        if score < 0.68:
            continue

        value_part = re.split(r"[:|=-]", line, maxsplit=1)
        raw_value = value_part[1] if len(value_part) > 1 else ""
        if not raw_value and index + 1 < len(lines):
            raw_value = lines[index + 1]

        value = cleaner(raw_value)
        if value:
            confidence = round(min(100.0, 55 + (score * 45)), 2)
            candidates.append(ExtractedField(value=value, confidence=confidence, source=line))

    if candidates:
        return max(candidates, key=lambda item: item.confidence)

    if fallback_pattern:
        joined_text = "\n".join(lines)
        match = re.search(fallback_pattern, joined_text, re.IGNORECASE | re.MULTILINE)
        if match:
            value = cleaner(match.group(0))
            if value:
                return ExtractedField(value=value, confidence=60.0, source="fallback-pattern")

    return None


def parse_id_card_text(text: str, ocr_confidence: float = 0.0) -> dict[str, ExtractedField]:
    """Heuristic parser for student ID style documents (falls back when AI fails).

    Returns a dict of canonical keys -> ExtractedField with reasonable confidences.
    """
    lines = [normalize_text(line) for line in text.splitlines() if normalize_text(line)]
    joined = "\n".join(lines)
    fields: dict[str, ExtractedField] = {}

    def setk(key: str, value: str, conf: float, src: str = "id-parser"):
        if not value:
            return
        value = clean_field_value(key, value)
        if not value:
            return
        fields[key] = ExtractedField(value=value, confidence=round(min(100.0, conf), 2), source=src)

    # Institution / card type
    if re.search(r"college of engineering|student identity card|student id", joined, re.I):
        inst = find_first_match(joined, r"^(.+COLLEGE OF ENGINEERING.+)$", flags=re.I | re.M)
        if not inst:
            inst = find_first_match(joined, r"^(.+COLLEGE.+)$", flags=re.I | re.M)
        if inst:
            setk("institution", inst, 90.0, "id-parser-institution")
        setk("card_type", "Student Identity Card", 95.0, "id-parser-cardtype")

    # Academic period
    ap = find_first_match(joined, r"(20\d{2}\s*[-–/]\s*20\d{2}|2023[-–]2027)")
    if ap:
        setk("academic_period", ap, 85.0, "id-parser-academic-period")

    # Name, father name, branch, student id, address, dob
    nm = find_first_match(joined, r"Name\s*[:\-]\s*(.+)", flags=re.I)
    if not nm:
        # sometimes just a line starting with uppercase words
        for line in lines[:8]:
            if re.fullmatch(r"[A-Z ]{3,}", line):
                nm = line
                break
    if nm:
        setk("name", nm, 92.0, "id-parser-name")

    fn = find_first_match(joined, r"F\.?\s*Name\s*[:\-]\s*(.+)|Father'?s?\s*name\s*[:\-]\s*(.+)", flags=re.I)
    if fn:
        setk("father_name", fn, 88.0, "id-parser-father")

    br = find_first_match(joined, r"Branch\s*[:\-]\s*(.+)", flags=re.I)
    if br:
        setk("branch", br, 86.0, "id-parser-branch")

    sid = None
    # look for a prominent uppercase alphanumeric token (roll/ID)
    for token in re.findall(r"\b[A-Z0-9]{6,12}\b", joined):
        if sum(c.isalpha() for c in token) >= 1 and sum(c.isdigit() for c in token) >= 1:
            sid = token
            break
    if not sid:
        sid = find_first_match(joined, r"Roll\s*[:\-]\s*([A-Z0-9-]+)", flags=re.I)
    if sid:
        setk("student_id", sid, 95.0, "id-parser-studentid")

    addr = find_first_match(joined, r"Address\s*[:\-]\s*(.+)", flags=re.I)
    if addr:
        setk("address", addr, 80.0, "id-parser-address")

    dob = find_first_match(joined, r"D\.?O\.?B\.?\s*[:\-]\s*(\d{1,2}[-/]\d{1,2}[-/]\d{2,4})", flags=re.I)
    if dob:
        # normalize to YYYY-MM-DD if possible
        try:
            d = datetime_from_string(dob)
            setk("date_of_birth", d, 90.0, "id-parser-dob")
        except Exception:
            setk("date_of_birth", dob, 70.0, "id-parser-dob")

    # phones
    phones = re.findall(r"(\d{10})", joined)
    if phones:
        # pick up to two numbers
        setk("phone", ",".join(phones[:2]), 92.0, "id-parser-phone")

    # signature: try to capture name after 'Student Signature' or 'Signature'
    sig = find_first_match(joined, r"Student\s+Signature\s*[:\-]?\s*(.+)$", flags=re.I | re.M)
    if not sig:
        # maybe on a later line
        m = re.search(r"Student\s+Signature\b.*\n\s*([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)", joined, re.I)
        if m:
            sig = m.group(1)
    if sig:
        setk("signature_text", sig, 80.0, "id-parser-signature")

    return fields


def find_first_match(text: str, pattern: str, flags=0) -> str | None:
    m = re.search(pattern, text, flags)
    if not m:
        return None
    # return the first non-empty capture group or the whole match
    for g in m.groups():
        if g:
            return g.strip()
    return m.group(0).strip()


def datetime_from_string(s: str) -> str:
    # Try common date formats and return ISO YYYY-MM-DD
    from datetime import datetime

    candidates = ["%d-%m-%Y", "%d/%m/%Y", "%d-%m-%y", "%d/%m/%y", "%Y-%m-%d"]
    for fmt in candidates:
        try:
            d = datetime.strptime(s, fmt)
            return d.strftime("%Y-%m-%d")
        except Exception:
            continue
    raise ValueError("unparseable date")


def extract_fields(text: str, ocr_confidence: float = 0.0) -> dict[str, ExtractedField]:
    lines = [normalize_text(line) for line in text.splitlines() if normalize_text(line)]

    extracted = {
        "first_name": extract_by_label(lines, "first_name", clean_name),
        "middle_name": extract_by_label(lines, "middle_name", clean_name),
        "last_name": extract_by_label(lines, "last_name", clean_name),
        "name": extract_by_label(lines, "name", clean_name),
        "gender": extract_by_label(lines, "gender", clean_plain_text),
        "date_of_birth": extract_by_label(lines, "date_of_birth", clean_plain_text),
        "address_line_1": extract_by_label(lines, "address_line_1", clean_plain_text),
        "address_line_2": extract_by_label(lines, "address_line_2", clean_plain_text),
        "city": extract_by_label(lines, "city", clean_plain_text),
        "state": extract_by_label(lines, "state", clean_plain_text),
        "pin_code": extract_by_label(lines, "pin_code", clean_pin_code, r"\b\d{5,6}\b"),
        "phone": extract_by_label(lines, "phone", clean_phone, r"(?:\+?\d[\d\s().-]{7,}\d)"),
        "email": extract_by_label(lines, "email", clean_email, r"[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}"),
    }

    result = {}
    for key, field in extracted.items():
        if field is None:
            continue

        combined_confidence = round((field.confidence * 0.65) + (ocr_confidence * 0.35), 2)
        result[key] = ExtractedField(
            value=field.value,
            confidence=combined_confidence,
            source=field.source,
        )

    return result


def extract_dynamic_fields_from_text(text: str, confidence: float = 70.0) -> dict[str, ExtractedField]:
    fields = {}
    lines = [normalize_text(line) for line in text.splitlines() if normalize_text(line)]

    email = find_first_email(text)
    phone = find_first_phone(text)
    name = guess_name_from_text(text, email, phone)

    if email:
        fields["email"] = ExtractedField(
            value=clean_email(email),
            confidence=95.0,
            source="resume-email-detection",
        )

    if phone:
        fields["phone"] = ExtractedField(
            value=clean_phone(phone),
            confidence=95.0,
            source="resume-phone-detection",
        )

    if name:
        fields["name"] = ExtractedField(
            value=clean_name(name),
            confidence=85.0,
            source="resume-name-detection",
        )

    for line in lines:
        if len(line) > 300:
            continue

        if re.search(r"@", line) and not re.search(r"\b(email|e-mail|mail|contact|phone|linkedin|github|leetcode|portfolio|website)\b", line, re.I):
            continue

        if re.search(r"\b(linkedin|github|leetcode|portfolio|website)\b", line, re.I):
            continue

        parts = re.split(r"\s*(?:[:=]|\t|\s\|\s)\s*", line, maxsplit=1)
        if len(parts) != 2:
            continue

        label, value = parts[0].strip(), parts[1].strip()
        if not label or not value or len(label) > 80 or should_skip_extracted_field(label, value):
            continue

        key = normalize_key(label)
        canonical_key = canonical_key_for_label(label)
        output_key = canonical_key or key
        value = clean_field_value(output_key, value)
        if not value:
            continue

        fields[output_key] = ExtractedField(
            value=value,
            confidence=confidence,
            source=label,
        )

    return fields


def merge_fields(page_results: list[dict]) -> dict[str, ExtractedField]:
    merged_fields: dict[str, ExtractedField] = {}

    for page in page_results:
        for key, field in page["fields"].items():
            if key not in merged_fields or field.confidence > merged_fields[key].confidence:
                merged_fields[key] = field

    return merged_fields


@app.post("/extract")
async def extract(file: UploadFile = File(...), page_screenshot: UploadFile | None = File(None)):
    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail="No file uploaded.")

    page_context = None
    if page_screenshot is not None:
        screenshot_bytes = await page_screenshot.read()
        if screenshot_bytes:
            try:
                screenshot_image = Image.open(io.BytesIO(screenshot_bytes))
                screenshot_image = screenshot_image.convert("RGB")
                page_context = build_page_context_from_screenshot(screenshot_image)
            except UnidentifiedImageError:
                page_context = {"error": "Unable to decode page screenshot."}

    extension = file_extension(file.filename)
    page_results = []

    if extension not in IMAGE_EXTENSIONS and extension != ".pdf" and not (file.content_type or "").startswith("image/"):
        text_upload = text_from_upload(file_bytes, file.filename)
        if text_upload is None:
            raise HTTPException(
                status_code=400,
                detail="Unsupported file type. Supported: images, PDF, TXT, CSV, JSON, HTML, MD, DOCX, XLSX, XLSM, PPTX.",
            )

        text, source_type = text_upload
        if not text.strip():
            raise HTTPException(status_code=400, detail="No readable text found in this file.")

        gemini_fields, gemini_status = extract_text_fields_with_gemini(text)
        local_fields = extract_dynamic_fields_from_text(text)
        page_results.append(
            {
                "page": 1,
                "raw_text": text,
                "ocr_confidence": 100.0,
                "extraction_method": "gemini" if gemini_fields else source_type,
                "gemini_status": gemini_status,
                "fields": gemini_fields or local_fields,
            }
        )

        return response_from_pages(page_results, page_context)

    images = images_from_upload(file_bytes, file.content_type, file.filename)

    for page_number, image in enumerate(images, start=1):
        gemini_fields, gemini_status = extract_fields_with_gemini(image)
        text, ocr_confidence = ocr_image(image)
        local_fields = extract_fields(text, ocr_confidence)
        id_fields = parse_id_card_text(text, ocr_confidence)

        # Merge id-parser fields into local fields, preferring higher confidence
        merged_local = {**local_fields}
        for k, v in id_fields.items():
            if k not in merged_local or v.confidence > merged_local[k].confidence:
                merged_local[k] = v

        if gemini_fields:
            final_fields = gemini_fields
            extraction_method = "gemini"
        else:
            final_fields = merged_local
            extraction_method = "local-ocr+id-parser" if id_fields else "local-ocr"

        page_results.append(
            {
                "page": page_number,
                "raw_text": text,
                "ocr_confidence": ocr_confidence,
                "extraction_method": extraction_method,
                "gemini_status": gemini_status,
                "fields": final_fields,
            }
        )

    return response_from_pages(page_results)


@app.post("/page_context")
async def page_context(page_screenshot: UploadFile = File(...)):
    """Return page context (hint labels and AI field label map) for a provided screenshot image.

    This allows the frontend to upload the page screenshot independently of the document
    extraction so both can proceed in parallel and be merged client-side.
    """
    screenshot_bytes = await page_screenshot.read()
    if not screenshot_bytes:
        raise HTTPException(status_code=400, detail="No screenshot uploaded.")

    try:
        screenshot_image = Image.open(io.BytesIO(screenshot_bytes))
        screenshot_image = screenshot_image.convert("RGB")
        page_ctx = build_page_context_from_screenshot(screenshot_image)
        return {"page_context": page_ctx}
    except UnidentifiedImageError:
        raise HTTPException(status_code=400, detail="Unable to decode page screenshot.")


@app.post("/map_fields")
async def map_fields(
    data: str = File(...),
    dom: str = File(...),
    page_screenshot: UploadFile | None = File(None),
    source_file: UploadFile | None = File(None),
    extracted_text: str | None = Form(None),
    fields_data: str | None = Form(None),
):
    """Map extracted data keys to DOM controls using optional page screenshot for AI assistance.

    - `data`: JSON string of extracted data `{key: value}`
    - `dom`: JSON string array of DOM control snapshots as produced by the content script
    - `page_screenshot`: optional image file to run AI page context extraction
    - `source_file`: optional uploaded source document to provide additional context for AI
    - `extracted_text`: optional extracted document text to help AI mapping
    """
    try:
        extracted = json.loads(data)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON in data field")

    try:
        dom_controls = json.loads(dom)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON in dom field")

    ai_label_map = {}
    screenshot_image = None
    mapping_method = "local"

    if page_screenshot is not None:
        screenshot_bytes = await page_screenshot.read()
        if screenshot_bytes:
            try:
                screenshot_image = Image.open(io.BytesIO(screenshot_bytes)).convert("RGB")
                page_ctx = build_page_context_from_screenshot(screenshot_image)
                ai_label_map = page_ctx.get("ai_field_label_map") or {}
            except UnidentifiedImageError:
                screenshot_image = None
                ai_label_map = {}

    source_filename = None
    if source_file is not None:
        source_filename = source_file.filename

    field_sources: dict[str, str] = {}
    if fields_data:
        try:
            parsed_fields = json.loads(fields_data)
            if isinstance(parsed_fields, dict):
                field_sources = {
                    key: item.get('source') if isinstance(item, dict) else ''
                    for key, item in parsed_fields.items()
                }
        except Exception:
            field_sources = {}

    ai_mappings = []
    if screenshot_image is not None:
        try:
            source_name = source_filename if source_filename else None
            extracted_text_value = extracted_text if extracted_text else None
            ai_mappings, ai_status = extract_field_mapping_with_gemini(
                screenshot_image,
                extracted,
                dom_controls,
                field_sources=field_sources,
                extracted_text=extracted_text_value,
                source_filename=source_name,
            )
            if ai_mappings:
                mapping_method = "ai"
            else:
                mapping_method = "local"
        except Exception:
            ai_mappings = []
            mapping_method = "local"

    if ai_mappings:
        return {"mappings": ai_mappings, "mapping_method": mapping_method}

    def control_text(ctrl: dict) -> str:
        parts = [ctrl.get('labelText') or '', ctrl.get('placeholder') or '', ctrl.get('nearestText') or '', ctrl.get('exactText') or '', ctrl.get('name') or '', ctrl.get('id') or '']
        return normalize_text(' '.join([p for p in parts if p]))

    mappings = []

    for key, value in extracted.items():
        best = None
        best_score = 0.0
        key_label = re.sub(r"_", " ", key)

        for ctrl in dom_controls:
            text = control_text(ctrl)
            if not text:
                continue

            s1 = similarity(text, key)
            s2 = similarity(text, key_label)
            s3 = 0.0
            # AI label contributions
            for ai_lbls in ai_label_map.get(key, []):
                s3 = max(s3, similarity(text, ai_lbls))

            score = max(s1, s2, s3)

            # bonus if id/name contains key tokens
            nid = (ctrl.get('id') or '') + ' ' + (ctrl.get('name') or '')
            if key.replace('_', '') and key.replace('_', '') in nid.replace('_', ''):
                score = max(score, 1.0)

            if score > best_score:
                best_score = score
                best = ctrl

        if best and best_score >= 0.50:
            selector = None
            if best.get('id'):
                selector = {"by": "id", "value": best.get('id')}
            elif best.get('name'):
                selector = {"by": "name", "value": best.get('name')}
            else:
                selector = {"by": "xpath", "value": best.get('xpath')}

            mappings.append({
                "key": key,
                "selector": selector,
                "score": round(best_score, 3),
                "control_text": control_text(best),
            })
        else:
            mappings.append({"key": key, "selector": None, "score": round(best_score, 3), "control_text": ""})

    return {"mappings": mappings, "mapping_method": mapping_method}